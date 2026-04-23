"""
sync_manager.py — Автосинхронізація Google Drive для EMET-бота.

Що робить:
- RAG sync: порівнює modifiedTime файлів у Drive з sync_state в PostgreSQL.
  При змінах — перебудовує ChromaDB-індекси у фоні (temp dir → атомарний swap).
  Сертифікати (certs) — тільки оновлення sync_state, без RAG-індексу (PDF скановані).
- Course sync: парсить Google Sheets з COURSE_FOLDER_ID → оновлює курси в PostgreSQL.

Формат курсу (Google Spreadsheet):
  Ім'я файлу = назва курсу
  Лист "теми":   стовпці: # | Назва теми | Зміст
  Лист "тести":  стовпці: Тема # | Питання | Варіант A | B | C | D | Правильна (A/B/C/D)

Конфігурація (.env):
  COURSE_FOLDER_ID   — ID папки Google Drive з курсами (якщо не задано — курси не синхронізуються)
  SYNC_INTERVAL_SEC  — інтервал перевірки в секундах (за замовчуванням 3600)
"""

import os
import io
import time
import shutil
import threading
import base64
import json as _json
import logging
import urllib.request
import urllib.parse
import db
import pandas as pd
from datetime import datetime

logger = logging.getLogger("emet_sync")


def _notify_admin(text: str):
    """Telegram-сповіщення адміну при критичних подіях sync. Не падає при помилці."""
    token = os.getenv("TELEGRAM_TOKEN")
    admin_id = os.getenv("ADMIN_ID")
    if not token or not admin_id:
        return
    try:
        data = urllib.parse.urlencode({"chat_id": admin_id, "text": text[:4000]}).encode()
        req = urllib.request.Request(
            f"https://api.telegram.org/bot{token}/sendMessage",
            data=data,
            method="POST"
        )
        urllib.request.urlopen(req, timeout=10)
    except Exception as e:
        logger.error(f"_notify_admin failed: {e}")

from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from dotenv import load_dotenv

from langchain_core.documents import Document
from langchain_chroma import Chroma
from langchain_openai import OpenAIEmbeddings
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter

load_dotenv()

OPENAI_KEY = os.getenv("OPENAI_API_KEY")
GEMINI_KEY = os.getenv("GEMINI_API_KEY")
COURSE_FOLDER_ID = os.getenv("COURSE_FOLDER_ID", "")
SYNC_INTERVAL_SEC = int(os.getenv("SYNC_INTERVAL_SEC", "3600"))

SERVICE_ACCOUNT_FILE = "credentials.json"
SCOPES = [
    "https://www.googleapis.com/auth/drive.readonly",
    "https://www.googleapis.com/auth/spreadsheets.readonly",
]

# Запобігає паралельній перебудові індексів (auto-sync + ручний sync з адмін-панелі)
_sync_lock = threading.Lock()

# Папки Google Drive → пути индексов + размер чанка
# label используется как суффикс в имени папки индекса: db_index_{label}_{provider}
# Примітка: certs навмисно виключені — пошук по сертифікатах йде через SQL (sync_state),
# бо PDF-файли сертифікатів скановані; RAG-індекс для них марний.
RAG_FOLDERS = {
    "kb_openai":    {"folder_id": "1RBXHGXOIc2kkSAw-LqzLaRqEE3Ix7L-m", "db": "data/db_index_kb_openai",    "chunk_size": 800,  "overlap": 150, "provider": "openai"},
    "kb_google":    {"folder_id": "1RBXHGXOIc2kkSAw-LqzLaRqEE3Ix7L-m", "db": "data/db_index_kb_google",    "chunk_size": 800,  "overlap": 150, "provider": "google"},
    "coach_openai": {"folder_id": "1KPPBurEoCV_wWzY5HxEtv_TrMI4qXfPa",  "db": "data/db_index_coach_openai", "chunk_size": 1200, "overlap": 300, "provider": "openai"},
    "coach_google": {"folder_id": "1KPPBurEoCV_wWzY5HxEtv_TrMI4qXfPa",  "db": "data/db_index_coach_google", "chunk_size": 1200, "overlap": 300, "provider": "google"},
}

# Папка сертифікатів — тільки для sync_state (SQL-пошук по іменах файлів)
CERTS_FOLDER_ID = "1ma-6CNO2FeHaicbRag7RvStkf5Rp1MyJ"

# ─── Авторизация ──────────────────────────────────────────────────────────────

def _get_google_credentials():
    """
    Завантажує Google Service Account credentials.
    Спочатку пробує env var GOOGLE_SERVICE_ACCOUNT_JSON (JSON-рядок або base64).
    Fallback: файл SERVICE_ACCOUNT_FILE (для локальної розробки).
    """
    creds_env = os.getenv("GOOGLE_SERVICE_ACCOUNT_JSON", "")
    if creds_env:
        try:
            info = _json.loads(creds_env)
        except Exception:
            info = _json.loads(base64.b64decode(creds_env).decode())
        return service_account.Credentials.from_service_account_info(info, scopes=SCOPES)
    return service_account.Credentials.from_service_account_file(SERVICE_ACCOUNT_FILE, scopes=SCOPES)


def get_services():
    creds = _get_google_credentials()
    drive = build("drive", "v3", credentials=creds)
    sheets = build("sheets", "v4", credentials=creds)
    return drive, sheets


# ─── Инициализация таблиц ─────────────────────────────────────────────────────

def init_sync_tables():
    """Проверка что таблицы sync_state и courses уже созданы в init_db() (main.py).
    Оставлена для обратной совместимости — вызывается из main.py."""
    pass  # Таблицы созданы в init_db() через db.get_connection()


# ─── Список файлов Drive с метаданными ───────────────────────────────────────

def list_files_with_meta(drive, folder_id):
    """Рекурсивно возвращает список файлов с полями id, name, mimeType, modifiedTime, webViewLink."""
    result = []
    page_token = None
    while True:
        try:
            resp = drive.files().list(
                q=f"'{folder_id}' in parents and trashed = false",
                fields="nextPageToken, files(id, name, mimeType, modifiedTime, webViewLink)",
                pageToken=page_token,
                includeItemsFromAllDrives=True,
                supportsAllDrives=True,
            ).execute()
            for item in resp.get("files", []):
                if item["mimeType"] == "application/vnd.google-apps.folder":
                    result.extend(list_files_with_meta(drive, item["id"]))
                else:
                    result.append(item)
            page_token = resp.get("nextPageToken")
            if not page_token:
                break
        except Exception as e:
            print(f"Ошибка доступа к папке {folder_id}: {e}")
            break
    return result


# ─── Извлечение текста из файлов (все поддерживаемые форматы) ────────────────

def extract_text(drive, file):
    mime = file["mimeType"]
    fid = file["id"]
    text = ""
    try:
        if mime == "application/pdf":
            buf = _download_bytes(drive, fid)
            reader = PdfReader(buf)
            text = "".join(p.extract_text() or "" for p in reader.pages)

        elif mime == "application/vnd.google-apps.document":
            text = drive.files().export_media(fileId=fid, mimeType="text/plain").execute().decode("utf-8")

        elif mime == "application/vnd.google-apps.spreadsheet":
            text = drive.files().export_media(fileId=fid, mimeType="text/csv").execute().decode("utf-8")

        elif mime == "application/vnd.google-apps.presentation":
            text = drive.files().export_media(fileId=fid, mimeType="text/plain").execute().decode("utf-8")

        elif mime == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            buf = _download_bytes(drive, fid)
            text = pd.read_excel(buf).to_csv(index=False)

        elif mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            buf = _download_bytes(drive, fid)
            prs = Presentation(buf)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

        elif mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            buf = _download_bytes(drive, fid)
            doc = DocxDocument(buf)
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())

        elif mime in ("text/plain", "text/csv"):
            text = drive.files().get_media(fileId=fid).execute().decode("utf-8", errors="replace")

    except Exception as e:
        print(f"Ошибка извлечения текста '{file['name']}': {e}")
    return text


MAX_FILE_BYTES = 50 * 1024 * 1024  # 50 MB

def _download_bytes(drive, file_id) -> io.BytesIO:
    buf = io.BytesIO()
    dl = MediaIoBaseDownload(buf, drive.files().get_media(fileId=file_id))
    done = False
    while not done:
        _, done = dl.next_chunk()
        if buf.tell() > MAX_FILE_BYTES:
            raise ValueError(f"Файл {file_id} перевищує ліміт {MAX_FILE_BYTES // 1024 // 1024} МБ, пропускаємо.")
    buf.seek(0)
    return buf


# ─── Split coach into products + competitors ─────────────────────────────────

COMP_PATTERNS = ["competitor", "competitir", "_master."]

def _split_coach_to_products_competitors():
    """After coach index is rebuilt, split it into products and competitors."""
    from langchain_chroma import Chroma as _Chroma
    from langchain_core.documents import Document as _Doc

    coach_dir = RAG_FOLDERS["coach_openai"]["db"]
    products_dir = "data/db_index_products_openai"
    competitors_dir = "data/db_index_competitors_openai"

    emb = OpenAIEmbeddings(model="text-embedding-3-small", openai_api_key=OPENAI_KEY)
    src = _Chroma(persist_directory=coach_dir, embedding_function=emb)
    data = src._collection.get(limit=5000, include=["metadatas", "documents"])

    if not data["ids"]:
        print("  [split] coach is empty, skipping")
        return

    products, competitors = [], []
    for meta, content in zip(data["metadatas"], data["documents"]):
        src_name = meta.get("source", "").lower()
        is_comp = any(p in src_name for p in COMP_PATTERNS)
        doc = _Doc(page_content=content, metadata=meta)
        (competitors if is_comp else products).append(doc)

    # Add LMS topics to products
    try:
        import db as _db
        splitter = RecursiveCharacterTextSplitter(chunk_size=1200, chunk_overlap=300)
        topics = _db.query_dict(
            "SELECT t.title, t.content, c.title as ct "
            "FROM topics t JOIN courses c ON c.id=t.course_id "
            "WHERE t.content IS NOT NULL AND length(trim(t.content)) > 50 "
            "ORDER BY c.id, t.order_num"
        )
        for t in topics:
            text = f"# {t['ct']}\n## {t['title']}\n\n{t['content']}"
            fn = f"[LMS] {t['ct']} -- {t['title']}"
            doc = _Doc(page_content=text, metadata={"source": fn, "url": "lms_course", "folder": "products"})
            products.extend(splitter.split_documents([doc]))
    except Exception as e:
        print(f"  [split] LMS add error: {e}")

    # Filter out heading-only / empty chunks (sync bug protection)
    def _has_real_content(d):
        # Видаляємо markdown-заголовки та пробіли — рахуємо реальні символи
        text = d.page_content
        for line in text.split("\n"):
            s = line.strip()
            if s and not s.startswith("#"):
                return len(text.strip()) >= 250  # мінімум 250 символів реального контенту
        return False

    products = [d for d in products if _has_real_content(d)]
    competitors = [d for d in competitors if _has_real_content(d)]

    # Enrich metadata with product_canonical for product-locked RAG retrieval
    def _detect_product_canonical(doc):
        """Визначає канонічний продукт з source filename + content."""
        src = (doc.metadata.get("source", "") or "").lower()
        text = (doc.page_content or "")[:500].lower()
        combined = src + " " + text

        # Vitaran варіанти — спочатку специфічні, потім generic
        if any(k in combined for k in ["whitening", "вайтенинг", "вайтенінг"]):
            return "HP Cell Vitaran Whitening"
        if any(k in combined for k in ["tox eye", "тохтай", "токс ай", "tox&face"]):
            return "HP Cell Vitaran Tox Eye"
        if any(k in combined for k in ["skin healer", "скін хілер"]):
            return "Vitaran Skin Healer"
        if any(k in combined for k in ["vitaran iii", "vitaran_iii", "vitaran ii", "vitaran_ii"]):
            return "HP Cell Vitaran iII"
        if any(k in combined for k in ["vitaran i ", "vitaran i\n", "vitaran_i", "hp cell vitaran i", "vitaran i.", "vitaran i,"]):
            return "HP Cell Vitaran i"
        if any(k in combined for k in ["vitaran", "вітаран", "витаран", "hp cell"]):
            return "Vitaran"

        if any(k in combined for k in ["ellans", "елансе", "ellanse"]):
            return "Ellansé"
        if any(k in combined for k in ["petaran", "петаран", "poly plla", "полі-l-молочна"]):
            return "Petaran"
        if any(k in combined for k in ["exoxe", "ексоксе", "экзосом"]):
            return "EXOXE"
        if "neuramis" in combined or "нейрамис" in combined or "нейраміс" in combined:
            return "Neuramis"
        if "neuronox" in combined or "нейронокс" in combined:
            return "Neuronox"
        if "iuse skin" in combined or "скінбустер" in combined or "skinbooster" in combined or "iuse_sb" in combined:
            return "IUSE SKINBOOSTER HA 20"
        if "iuse hair" in combined:
            return "IUSE HAIR REGROWTH"
        if "iuse collagen" in combined:
            return "IUSE Collagen"
        if "esse" in combined or "ессе" in combined:
            return "ESSE"
        if "magnox" in combined or "магнокс" in combined:
            return "Magnox"
        if "iuse" in combined:
            return "IUSE"
        return None

    def _detect_scope(doc):
        """Визначає рівень специфічності чанка: line / product / ingredient / protocol.
        line — загальні характеристики лінії (часто одночасна згадка >=2 продуктів того ж бренду без фокусу)
        product — конкретний продукт (одна назва домінує)
        ingredient — про окремий компонент (PLLA, PDRN, PCL...)
        protocol — про процедуру / схему / комбо
        """
        src = (doc.metadata.get("source", "") or "").lower()
        text = (doc.page_content or "").lower()
        # protocol — за source або тексту
        if any(k in src for k in ["комбін", "протокол", "combo", "protokol"]):
            return "protocol"
        if any(k in text[:200] for k in ["протокол", "розведення", "схема процедур", "техніка"]):
            return "protocol"
        # ingredient — пояснення про компонент окремо
        if any(k in text[:300] for k in [" plla ", "поліl-молочна", "поликапролак", "пдрн", " pdrn ",
                                          " pcl ", "поликапролактон", "гіалуронова кислот", "hyaluronic"]):
            if not any(p in text[:100] for p in ["petaran", "петаран", "ellans", "елансе", "vitaran", "вітаран",
                                                   "neuramis", "нейрамис", "iuse"]):
                return "ingredient"
        # line — згадка >=2 продуктів того самого бренду / lineup keywords
        line_markers = [
            ("esse", ["sensitive", "sensitive plus", "core", "professional", "лінійка esse", "лінія esse",
                      "пробіотична космецевтика", "лінійки", "асортимент"]),
            ("vitaran", ["лінійка vitaran", "лінія vitaran", "all variants", "усі варіанти"]),
            ("iuse", ["лінійка iuse", "лінія iuse", "skinbooster і hair", "колаген і hair"]),
        ]
        for brand, markers in line_markers:
            if brand in text[:400] and any(m in text[:400] for m in markers):
                return "line"
        # default
        return "product"

    for d in products + competitors:
        canonical = _detect_product_canonical(d)
        if canonical:
            d.metadata["product_canonical"] = canonical
        d.metadata["scope"] = _detect_scope(d)

    # Лог розподілу по продуктам + scope
    from collections import Counter
    prod_dist = Counter(d.metadata.get("product_canonical", "UNKNOWN") for d in products)
    scope_dist = Counter(d.metadata.get("scope", "?") for d in products + competitors)
    print(f"  [split] products by canonical: {dict(prod_dist.most_common())}")
    print(f"  [split] scope distribution: {dict(scope_dist.most_common())}")

    # Rebuild both indices
    counts = {}
    for path, docs, label in [(products_dir, products, "products"), (competitors_dir, competitors, "competitors")]:
        shutil.rmtree(path, ignore_errors=True)
        vdb = _Chroma(persist_directory=path, embedding_function=emb)
        BATCH = 50
        for i in range(0, len(docs), BATCH):
            _batch_to_chroma_simple(docs[i:i+BATCH], emb, vdb)
        cnt = vdb._collection.count()
        counts[label] = cnt
        print(f"  [split] {label}: {cnt} chunks")

    # Sanity check: обидва індекси не повинні бути <100 чанків — це ознака провалу embedding
    SPLIT_MIN_CHUNKS = 100
    if counts.get("products", 0) < SPLIT_MIN_CHUNKS or counts.get("competitors", 0) < SPLIT_MIN_CHUNKS:
        msg = (f"SPLIT FAILED: products={counts.get('products', 0)} chunks, "
               f"competitors={counts.get('competitors', 0)} chunks (min {SPLIT_MIN_CHUNKS}). "
               f"Likely OpenAI rate limit or API failure during embedding.")
        logger.error(msg)
        _notify_admin(f"⚠️ EMET sync: {msg}")
        raise RuntimeError(msg)


def _batch_to_chroma_simple(docs, emb, vdb, rate_limit_sleep=10, max_retries=10):
    """Batch add з retry на rate-limit. Кидає виняток коли retries вичерпані."""
    rate_limit_keys = ["429", "RateLimitError", "rate_limit", "RESOURCE_EXHAUSTED"]
    retries = 0
    while True:
        try:
            vdb.add_documents(docs)
            return
        except Exception as e:
            err = str(e)
            if any(k in err for k in rate_limit_keys) and retries < max_retries:
                retries += 1
                logger.warning(f"  [split batch] rate limit, sleep {rate_limit_sleep}s (retry {retries}/{max_retries})")
                time.sleep(rate_limit_sleep)
            else:
                logger.error(f"  [split batch] failed after {retries} retries: {err[:200]}")
                raise


# ─── Построение индексов ──────────────────────────────────────────────────────

def _build_index(drive, files, cfg, folder_label):
    """Строит один индекс (OpenAI или Google) для переданного набора файлов."""
    docs = _files_to_documents(drive, files, folder_label)
    if not docs:
        return
    chunks = RecursiveCharacterTextSplitter(
        chunk_size=cfg["chunk_size"], chunk_overlap=cfg["overlap"]
    ).split_documents(docs)
    target_dir = cfg["db"]
    print(f"  [{cfg['provider']}] {len(chunks)} чанків → {target_dir}")
    if cfg["provider"] == "openai":
        emb = OpenAIEmbeddings(model="text-embedding-3-small", openai_api_key=OPENAI_KEY)
        _batch_to_chroma(chunks, emb, target_dir, rate_limit_sleep=10, rate_limit_keywords=["429", "RateLimitError"])
    else:
        emb = GoogleGenerativeAIEmbeddings(model="models/gemini-embedding-001", google_api_key=GEMINI_KEY)
        _batch_to_chroma(chunks, emb, target_dir, rate_limit_sleep=30, rate_limit_keywords=["429", "RESOURCE_EXHAUSTED", "503", "UNAVAILABLE"])


def _files_to_documents(drive, files, folder_label=""):
    documents = []
    for f in files:
        text = extract_text(drive, f)
        if text:
            name_lower = f["name"].lower()
            category = "combo" if any(w in name_lower for w in ["протокол", "комбін", "combo"]) else "general"
            # Для сертифікатів — тільки назва файлу (PDF-тіло часто скановане/garbled)
            if folder_label == "certs":
                page_content = f"Документ: {f['name']}"
            else:
                page_content = f"Документ: {f['name']}\n\n{text}"
            documents.append(Document(
                page_content=page_content,
                metadata={
                    "source":    f["name"],
                    "url":       f.get("webViewLink", ""),
                    "file_id":   f["id"],
                    "category":  category,
                    "folder":    folder_label,
                }
            ))
    return documents


MAX_RATE_LIMIT_RETRIES = 10

def _batch_to_chroma(chunks, embeddings, persist_dir, rate_limit_sleep, rate_limit_keywords):
    db = None
    for i in range(0, len(chunks), 50):
        batch = chunks[i:i + 50]
        retries = 0
        success = False
        while not success:
            try:
                if db is None:
                    db = Chroma.from_documents(batch, embedding=embeddings, persist_directory=persist_dir)
                else:
                    db.add_documents(batch)
                success = True
                time.sleep(0.5)
            except Exception as e:
                err = str(e)
                if any(k in err for k in rate_limit_keywords):
                    retries += 1
                    if retries > MAX_RATE_LIMIT_RETRIES:
                        raise RuntimeError(
                            f"Rate limit перевищено {MAX_RATE_LIMIT_RETRIES} разів поспіль для {persist_dir}. Зупиняємо."
                        ) from e
                    print(f"Rate limit / недоступність. Спроба {retries}/{MAX_RATE_LIMIT_RETRIES}. Сплю {rate_limit_sleep} сек...")
                    time.sleep(rate_limit_sleep)
                else:
                    raise


# ─── RAG синхронизация ────────────────────────────────────────────────────────

def sync_rag_indexes():
    """
    Перевіряє зміни в кожній RAG-папці окремо.
    Пересобирає тільки ті індекси, де є зміни (атомарний swap).
    Для certs — тільки оновлює sync_state (SQL-пошук), без побудови RAG-індексу.
    Повертає (all_changed_names, changed_by_category).
    """
    drive, _ = get_services()

    # Поточний стан БД
    indexed = {r[0]: r[1] for r in db.query("SELECT file_id, modified_time FROM sync_state")}

    all_changed_names = []
    changed_by_category = {}  # {"kb": 5, "coach": 3, "certs": 2}

    # ── 1. RAG-індекси (kb + coach) ──────────────────────────────────────────
    # Групуємо конфіги по folder_id щоб не сканувати одну папку двічі
    by_folder = {}
    for label, cfg in RAG_FOLDERS.items():
        fid = cfg["folder_id"]
        if fid not in by_folder:
            by_folder[fid] = {"files": None, "labels": []}
        by_folder[fid]["labels"].append(label)

    for folder_id, info in by_folder.items():
        files = list_files_with_meta(drive, folder_id)
        info["files"] = files
        changed = [f for f in files if indexed.get(f["id"]) != f["modifiedTime"]]

        if not changed:
            print(f"RAG sync: немає змін у папці {folder_id} ({len(files)} файлів)")
            continue

        folder_label = info["labels"][0].split("_")[0]  # "kb", "coach"
        print(f"RAG sync: {len(changed)} змін у [{folder_label}]. Перебудовую індекси...")

        for label in info["labels"]:
            cfg = RAG_FOLDERS[label]
            tmp_dir = cfg["db"] + "_building"
            shutil.rmtree(tmp_dir, ignore_errors=True)
            try:
                _build_index(drive, files, {**cfg, "db": tmp_dir}, folder_label)
                # Атомарний swap
                old_dir = cfg["db"] + "_old"
                shutil.rmtree(old_dir, ignore_errors=True)
                if os.path.exists(cfg["db"]):
                    os.rename(cfg["db"], old_dir)
                os.rename(tmp_dir, cfg["db"])
                shutil.rmtree(old_dir, ignore_errors=True)
                print(f"  [{label}] індекс оновлено: {cfg['db']}")
            except Exception as e:
                print(f"  [{label}] помилка побудови: {e}")
                shutil.rmtree(tmp_dir, ignore_errors=True)

        # Оновлюємо sync_state
        db.executemany(
            "INSERT INTO sync_state (file_id, file_name, modified_time, indexed_at, folder_label) VALUES (%s,%s,%s,%s,%s) "
            "ON CONFLICT (file_id) DO UPDATE SET file_name=EXCLUDED.file_name, "
            "modified_time=EXCLUDED.modified_time, indexed_at=EXCLUDED.indexed_at, folder_label=EXCLUDED.folder_label",
            [(f["id"], f["name"], f["modifiedTime"], datetime.now().isoformat(), folder_label) for f in files]
        )
        all_changed_names.extend(f["name"] for f in changed)
        changed_by_category[folder_label] = len(changed)

        # After coach rebuild — auto-split into products + competitors indices
        if folder_label == "coach":
            try:
                _split_coach_to_products_competitors()
                try:
                    from tests.test_knowledge_integrity import run_integrity_check
                    passed, report = run_integrity_check(verbose=False)
                    if passed:
                        print("  [integrity] OK — zero data loss")
                    else:
                        logger.error(f"  [integrity] FAILED!\n{report}")
                        _notify_admin(f"⚠️ EMET sync: integrity check FAILED after split.\n{report[:3500]}")
                except Exception as e:
                    logger.error(f"  [integrity] check error: {e}")
                    _notify_admin(f"⚠️ EMET sync: integrity check crashed: {e}")
            except Exception as e:
                logger.error(f"  [split] помилка: {e}")
                _notify_admin(f"⚠️ EMET sync: split failed: {e}")

    # ── 2. Certs — тільки sync_state, без RAG-індексу ────────────────────────
    if CERTS_FOLDER_ID:
        certs_files = list_files_with_meta(drive, CERTS_FOLDER_ID)
        certs_changed = [f for f in certs_files if indexed.get(f["id"]) != f["modifiedTime"]]
        if certs_changed:
            print(f"Certs sync: {len(certs_changed)} змін. Оновлюю sync_state...")
            db.executemany(
                "INSERT INTO sync_state (file_id, file_name, modified_time, indexed_at, folder_label) VALUES (%s,%s,%s,%s,%s) "
                "ON CONFLICT (file_id) DO UPDATE SET file_name=EXCLUDED.file_name, "
                "modified_time=EXCLUDED.modified_time, indexed_at=EXCLUDED.indexed_at, folder_label=EXCLUDED.folder_label",
                [(f["id"], f["name"], f["modifiedTime"], datetime.now().isoformat(), "certs") for f in certs_files]
            )
            all_changed_names.extend(f["name"] for f in certs_changed)
            changed_by_category["certs"] = len(certs_changed)
        else:
            print(f"Certs sync: немає змін ({len(certs_files)} файлів)")

    if all_changed_names:
        print(f"RAG sync завершено. Змінено: {all_changed_names}")
    return all_changed_names, changed_by_category


# ─── Синхронизация курсов из Google Sheets ───────────────────────────────────

def _parse_course_spreadsheet(sheets_svc, spreadsheet_id, course_title):
    """
    Читает листы "теми" и "тести" из спредшита.
    Возвращает структуру курса: {title, topics: [{title, content, questions}]}
    """
    def read_sheet(sheet_name):
        try:
            result = sheets_svc.spreadsheets().values().get(
                spreadsheetId=spreadsheet_id,
                range=f"'{sheet_name}'!A1:Z1000"
            ).execute()
            return result.get("values", [])
        except Exception as e:
            print(f"Лист '{sheet_name}' не найден в '{course_title}': {e}")
            return []

    # Парсим темы: # | Назва теми | Зміст
    topics = {}
    for row in read_sheet("теми")[1:]:  # пропускаем заголовок
        if len(row) < 3:
            continue
        try:
            order = int(str(row[0]).strip())
            topics[order] = {
                "title": str(row[1]).strip(),
                "content": str(row[2]).strip(),
                "questions": []
            }
        except ValueError:
            continue

    # Парсим тесты: Тема # | Питання | A | B | C | D | Правильна
    letter_to_idx = {"A": 0, "B": 1, "C": 2, "D": 3}
    for row in read_sheet("тести")[1:]:
        if len(row) < 7:
            continue
        try:
            topic_num = int(str(row[0]).strip())
        except ValueError:
            continue
        if topic_num not in topics:
            continue

        q_text = str(row[1]).strip()
        opts_raw = [str(row[i]).strip() if i < len(row) else "" for i in range(2, 6)]
        correct_letter = str(row[6]).strip().upper() if len(row) > 6 else ""
        correct_idx = letter_to_idx.get(correct_letter, -1)

        options = [(text, i == correct_idx) for i, text in enumerate(opts_raw) if text]
        if q_text and options:
            topics[topic_num]["questions"].append({"text": q_text, "options": options})

    return {"title": course_title, "topics": [topics[k] for k in sorted(topics.keys())]}


def _upsert_course(course_data, drive_file_id, drive_modified):
    """Удаляет старую версию курса и вставляет новую."""
    with db.get_connection() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT id FROM courses WHERE drive_file_id=%s", (drive_file_id,))
            existing = cur.fetchone()

            if existing:
                old_id = existing[0]
                cur.execute("SELECT id FROM topics WHERE course_id=%s", (old_id,))
                topic_ids = [r[0] for r in cur.fetchall()]
                for tid in topic_ids:
                    cur.execute("SELECT id FROM questions WHERE topic_id=%s", (tid,))
                    q_ids = [r[0] for r in cur.fetchall()]
                    for qid in q_ids:
                        cur.execute("DELETE FROM answer_options WHERE question_id=%s", (qid,))
                    cur.execute("DELETE FROM questions WHERE topic_id=%s", (tid,))
                cur.execute("DELETE FROM topics WHERE course_id=%s", (old_id,))
                cur.execute("DELETE FROM courses WHERE id=%s", (old_id,))

            cur.execute(
                "INSERT INTO courses (title, description, drive_file_id, drive_modified, created_at) "
                "VALUES (%s,%s,%s,%s,%s) RETURNING id",
                (course_data["title"], "", drive_file_id, drive_modified, datetime.now().isoformat())
            )
            course_id = cur.fetchone()[0]

            for order, topic in enumerate(course_data["topics"], 1):
                cur.execute(
                    "INSERT INTO topics (course_id, order_num, title, content) VALUES (%s,%s,%s,%s) RETURNING id",
                    (course_id, order, topic["title"], topic["content"])
                )
                topic_id = cur.fetchone()[0]
                for q in topic["questions"]:
                    cur.execute(
                        "INSERT INTO questions (topic_id, text) VALUES (%s,%s) RETURNING id",
                        (topic_id, q["text"])
                    )
                    q_id = cur.fetchone()[0]
                    for opt_text, is_correct in q["options"]:
                        cur.execute(
                            "INSERT INTO answer_options (question_id, text, is_correct) VALUES (%s,%s,%s)",
                            (q_id, opt_text, int(is_correct))
                        )


def sync_courses():
    """
    Проверяет Google Sheets в COURSE_FOLDER_ID на изменения.
    Обновляет курсы в SQLite.
    Возвращает список обновлённых курсов.
    """
    if not COURSE_FOLDER_ID:
        return []

    drive, sheets_svc = get_services()

    spreadsheets = [
        f for f in list_files_with_meta(drive, COURSE_FOLDER_ID)
        if f["mimeType"] == "application/vnd.google-apps.spreadsheet"
    ]

    existing = {
        r[0]: r[1]
        for r in db.query("SELECT drive_file_id, drive_modified FROM courses WHERE drive_file_id IS NOT NULL")
    }

    updated = []
    for f in spreadsheets:
        if existing.get(f["id"]) == f["modifiedTime"]:
            continue  # Не изменился

        try:
            course_data = _parse_course_spreadsheet(sheets_svc, f["id"], f["name"])
            if not course_data["topics"]:
                print(f"Курс '{f['name']}': лист 'теми' пустой или не найден, пропускаю")
                continue
            _upsert_course(course_data, f["id"], f["modifiedTime"])
            updated.append(f["name"])
            total_q = sum(len(t["questions"]) for t in course_data["topics"])
            print(f"Курс обновлён: '{f['name']}' ({len(course_data['topics'])} тем, {total_q} вопросов)")
        except Exception as e:
            print(f"Ошибка парсинга курса '{f['name']}': {e}")

    return updated


# ─── Главная функция синхронизации ───────────────────────────────────────────

def run_sync():
    """
    Запускает полную синхронизацию: RAG-индексы + курсы.
    Возвращает {'rag_updated': [...], 'courses_updated': [...], 'error': None | str}
    Если синхронизация уже выполняется — возвращает error='sync_in_progress' немедленно.
    """
    if not _sync_lock.acquire(blocking=False):
        print("[sync] Синхронізація вже виконується, пропускаємо.")
        return {"rag_updated": [], "rag_by_category": {}, "courses_updated": [], "error": "sync_in_progress"}
    try:
        result = {"rag_updated": [], "rag_by_category": {}, "courses_updated": [], "error": None}
        result["rag_updated"], result["rag_by_category"] = sync_rag_indexes()
        result["courses_updated"] = sync_courses()
    except Exception as e:
        result["error"] = str(e)
        print(f"Ошибка sync: {e}")
    finally:
        _sync_lock.release()
    return result