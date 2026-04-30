"""tools/smart_import.py — структурні імпортери для xlsx / pptx / docx.

Замінює "сирий текст → split 1200 chars" на handlers що поважають структуру файлу:
- xlsx: рядок = chunk (картка), не CSV-склейки
- pptx: слайд = chunk (заголовок + bullets разом)
- docx: split по H1/H2 заголовках, fallback на 1200 chars

Кожен chunk отримує rich metadata:
- source: ім'я файлу
- chunk_type: 'xlsx_row' | 'pptx_slide' | 'docx_section' | 'raw'
- product_canonical: спроба визначити продукт
- product_subline: для ESSE — лінія (Core/Sensitive/Plus)

Використання — імпортуй в sync_manager:
    from tools.smart_import import smart_extract_documents
    docs = smart_extract_documents(file_path, file_bytes, mime_type)
"""
from __future__ import annotations

import io
import re
from typing import Iterator
from langchain_core.documents import Document


# ============================================================
# XLSX — кожний рядок = одна картка
# ============================================================

def extract_xlsx_rows(file_bytes: bytes, source_name: str) -> Iterator[Document]:
    """xlsx → один Document на рядок. Зберігає всі колонки як structured text."""
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(file_bytes), data_only=True)
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        if ws.max_row < 2:
            continue
        # Headers — рядок 1, нормалізовані
        headers = []
        for c in range(1, ws.max_column + 1):
            h = ws.cell(row=1, column=c).value
            headers.append(str(h).strip() if h else f"col{c}")
        # Data rows
        for r in range(2, ws.max_row + 1):
            row_data = {}
            for c, h in enumerate(headers, 1):
                v = ws.cell(row=r, column=c).value
                if v is not None and str(v).strip():
                    row_data[h] = str(v).strip()
            if not row_data:
                continue
            # Перший непорожній field = "name" продукту (часто Найменування / Назва)
            name_keys = ["Найменування", "Назва", "Name", "Product", "Продукт"]
            name = next((row_data.get(k) for k in name_keys if k in row_data), None)
            if not name:
                # Беремо першу непорожню колонку
                name = next(iter(row_data.values()))
            # Формуємо markdown
            md_lines = [f"# {name[:120]}"]
            md_lines.append(f"_(з документа: {source_name}, аркуш: {sheet_name})_\n")
            for k, v in row_data.items():
                if k in name_keys:
                    continue
                md_lines.append(f"**{k}:** {v}")
            content = "\n".join(md_lines)
            # Detect product_canonical (source_name has priority over content)
            product_canonical = _detect_product_canonical_from_text(name + " " + content, source_name)
            product_subline = _detect_subline(sheet_name, content)
            yield Document(
                page_content=content,
                metadata={
                    "source": source_name,
                    "sheet": sheet_name,
                    "row": r,
                    "chunk_type": "xlsx_row",
                    "product_canonical": product_canonical,
                    "product_subline": product_subline,
                    "url": "xlsx_structured",
                    "folder": "products",
                }
            )


# ============================================================
# PPTX — кожний слайд = chunk (заголовок + bullets)
# ============================================================

def extract_pptx_slides(file_bytes: bytes, source_name: str) -> Iterator[Document]:
    """pptx → один Document на слайд. Зберігає заголовок + усі text-frames."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    for slide_num, slide in enumerate(prs.slides, 1):
        # Збираємо всі тексти зі слайду
        title = ""
        body_texts = []
        for shape in slide.shapes:
            if not hasattr(shape, "text") or not shape.text.strip():
                continue
            text = shape.text.strip()
            # Heuristic: title placeholder = найбільший шрифт або перший
            if shape.has_text_frame and (
                getattr(shape, "is_placeholder", False) and
                getattr(shape.placeholder_format, "idx", -1) == 0
            ):
                title = text
            else:
                body_texts.append(text)
        if not title and body_texts:
            title = body_texts.pop(0).split("\n")[0][:100]
        if not title and not body_texts:
            continue
        # Формуємо markdown
        md = f"# Слайд {slide_num}: {title or '(без заголовка)'}\n\n"
        md += f"_(з презентації: {source_name})_\n\n"
        for txt in body_texts:
            # Зберігаємо буллети — кожен рядок як bullet
            for line in txt.split("\n"):
                line = line.strip()
                if line:
                    md += f"- {line}\n"
        full_text = title + " " + " ".join(body_texts)
        yield Document(
            page_content=md,
            metadata={
                "source": source_name,
                "slide": slide_num,
                "chunk_type": "pptx_slide",
                "product_canonical": _detect_product_canonical_from_text(source_name + " " + full_text, source_name),
                "product_subline": _detect_subline("", full_text),
                "url": "pptx_structured",
                "folder": "products",
            }
        )


# ============================================================
# DOCX — split по H1/H2 заголовках, fallback на 1200 chars
# ============================================================

def extract_docx_sections(file_bytes: bytes, source_name: str) -> Iterator[Document]:
    """docx → один Document на секцію (H1/H2). Якщо немає заголовків — fallback на параграфи."""
    from docx import Document as DocxDocument
    doc = DocxDocument(io.BytesIO(file_bytes))
    sections = []  # list of (heading, list of paragraphs)
    current_heading = "Документ"
    current_paras = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = (para.style.name or "").lower() if para.style else ""
        if "heading 1" in style or "heading 2" in style or "title" in style:
            # Зберігаємо попередню секцію
            if current_paras:
                sections.append((current_heading, current_paras))
            current_heading = text[:120]
            current_paras = []
        else:
            current_paras.append(text)
    if current_paras:
        sections.append((current_heading, current_paras))
    # Якщо нема структури — fallback на 1200-char split
    if len(sections) <= 1 and current_paras:
        full_text = "\n".join(current_paras)
        # Розрізаємо по 1200 chars + 300 overlap (як було раніше)
        chunk_size = 1200
        overlap = 300
        for i in range(0, len(full_text), chunk_size - overlap):
            chunk_text = full_text[i:i + chunk_size]
            if not chunk_text.strip():
                continue
            yield Document(
                page_content=f"# {source_name}\n\n{chunk_text}",
                metadata={
                    "source": source_name,
                    "chunk_type": "docx_chunked",
                    "product_canonical": _detect_product_canonical_from_text(source_name + " " + chunk_text, source_name),
                    "product_subline": "",
                    "url": "docx_fallback",
                    "folder": "products",
                }
            )
        return
    # Нормальний шлях — секції за заголовками
    for heading, paras in sections:
        content_text = "\n\n".join(paras)
        if len(content_text.strip()) < 30:
            continue
        md = f"# {heading}\n_(з документа: {source_name})_\n\n{content_text}"
        yield Document(
            page_content=md,
            metadata={
                "source": source_name,
                "section": heading,
                "chunk_type": "docx_section",
                "product_canonical": _detect_product_canonical_from_text(source_name + " " + heading + " " + content_text, source_name),
                "product_subline": _detect_subline(heading, content_text),
                "url": "docx_structured",
                "folder": "products",
            }
        )


# ============================================================
# Helpers
# ============================================================

def _detect_product_canonical_from_text(text: str, source_name: str = "") -> str:
    """Мапінг ключових слів → канонічна назва EMET-продукту.

    source_name (filename) має пріоритет над content — для конкурентних docs:
    "Neuronox_Competitors_MASTER.docx" → Neuronox, навіть якщо в тексті порівнюють з Neuramis.
    """
    src = source_name.lower()
    t = text.lower()

    # ── STEP 1: Strong filename signals (приоритет для product-specific файлів) ──
    if "neuronox" in src or "нейронокс" in src:
        return "Neuronox"
    if any(k in src for k in ["exoxe", "ехохе", "ексоксе", "_exoxe_", "exoxe_"]):
        return "EXOXE"
    if "neuramis" in src or "нейрамис" in src or "нейраміс" in src:
        return "Neuramis"
    if "magnox" in src or "магнокс" in src:
        return "Magnox"
    if any(k in src for k in ["whitening", "вайтенинг", "вайтенінг"]):
        return "HP Cell Vitaran Whitening"
    if any(k in src for k in ["tox eye", "тохтай", "токс ай", "vitaran tox", "_tox_", "vitaran_tox"]):
        return "HP Cell Vitaran Tox Eye"
    if any(k in src for k in ["skin healer", "dual serum", "vitaran exosome",
                               "azulene", "sleeping cream", "wrapping serum"]):
        return "Vitaran Skin Healer"
    if any(k in src for k in ["petaran", "петаран"]):
        return "Petaran"
    if any(k in src for k in ["ellans", "ellanse", "елансе"]):
        return "Ellansé"
    if "iuse hair" in src or "iuse_hair" in src:
        return "IUSE HAIR REGROWTH"
    if "iuse collagen" in src or "iuse_collagen" in src:
        return "IUSE Collagen"
    if "iuse skin" in src or "skinbooster" in src or "skin booster" in src or "iuse_sb" in src or "впв skin" in src:
        return "IUSE SKINBOOSTER HA 20"
    if "esse" in src or "ессе" in src:
        return "ESSE"
    if any(k in src for k in ["vitaran iii", "vitaran_iii", "vitaran ii", "vitaran_ii"]):
        return "HP Cell Vitaran i"
    if any(k in src for k in ["vitaran i ", "vitaran_i", "hp cell vitaran"]):
        return "HP Cell Vitaran i"

    # ── STEP 2: Content matching (для файлів без явних маркерів в назві) ──
    if any(k in t for k in ["whitening", "вайтенинг", "вайтенінг"]):
        return "HP Cell Vitaran Whitening"
    if any(k in t for k in ["tox eye", "тохтай", "токс ай"]):
        return "HP Cell Vitaran Tox Eye"
    if any(k in t for k in ["skin healer", "vitaran exosome", "dual serum",
                             "azulene", "sleeping cream", "wrapping serum"]):
        return "Vitaran Skin Healer"
    if any(k in t for k in ["vitaran iii", "vitaran_iii", "vitaran ii", "vitaran_ii"]):
        return "HP Cell Vitaran i"
    if any(k in t for k in ["vitaran i ", "vitaran i\n", "vitaran_i", "hp cell vitaran"]):
        return "HP Cell Vitaran i"
    if "vitaran" in t or "вітаран" in t:
        return "Vitaran"
    if any(k in t for k in ["ellans", "елансе", "ellanse"]):
        return "Ellansé"
    if any(k in t for k in ["petaran", "петаран"]):
        return "Petaran"
    if any(k in t for k in ["exoxe", "ехохе", "ексоксе", "экзосом"]):
        return "EXOXE"
    if "neuronox" in t or "нейронокс" in t:
        return "Neuronox"
    if "neuramis" in t or "нейрамис" in t or "нейраміс" in t:
        return "Neuramis"
    if "iuse skin" in t or "скінбустер" in t or "skinbooster" in t or "skin booster" in t:
        return "IUSE SKINBOOSTER HA 20"
    if "iuse hair" in t or "iuse_hair" in t:
        return "IUSE HAIR REGROWTH"
    if "iuse collagen" in t:
        return "IUSE Collagen"
    if "esse" in t or "ессе" in t:
        return "ESSE"
    if "magnox" in t or "магнокс" in t:
        return "Magnox"
    return ""


def _detect_subline(sheet_name: str, content: str) -> str:
    """Для ESSE визначає лінію (Core/Sensitive/Plus/тощо).

    Note: cyrillic Е (U+0415) часто міксується з ASCII E в назвах ('Лінійка ЕSSE CORE').
    Тому нормалізуємо: cyrillic Е→E, е→e перед matching.
    """
    combined = (sheet_name + " " + content[:300]).lower()
    # Replace Ukrainian cyrillic e/Е with ASCII counterparts (e.g. "еsse" → "esse")
    combined = combined.replace("е", "e").replace("Е", "E")  # Cyrillic e/Е → ASCII

    if "sensitive" in combined or "сeнсітів" in combined or "сeнситив" in combined:
        return "Esse Sensitive"
    if "esse plus" in combined or "esse+" in combined or "лінійка esse plus" in combined:
        return "Esse Plus"
    if "professional" in combined or "профeсійн" in combined:
        return "Esse Professional"
    if "esse core" in combined or "лінійка esse core" in combined:
        return "Esse Core"
    if "консилeр" in combined or "concealer" in combined:
        return "Esse Concealers"
    if "сонцeзахис" in combined or "foundation" in combined or "тональн" in combined:
        return "Esse Sun Protection"
    if "набор" in combined or "набір" in combined or "set" in combined:
        return "Esse Sets"
    return ""


# ============================================================
# Public dispatcher
# ============================================================

def smart_extract_documents(file_bytes: bytes, source_name: str, mime_type: str) -> list[Document]:
    """Dispatcher — викликає правильний extractor за mime/extension."""
    src_lower = source_name.lower()
    try:
        if mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet" or src_lower.endswith(".xlsx"):
            return list(extract_xlsx_rows(file_bytes, source_name))
        if mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation" or src_lower.endswith(".pptx"):
            return list(extract_pptx_slides(file_bytes, source_name))
        if mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" or src_lower.endswith(".docx"):
            return list(extract_docx_sections(file_bytes, source_name))
    except Exception as e:
        # Fail-soft: повертаємо пусто щоб старий fallback в sync_manager міг спробувати
        import logging
        logging.getLogger("emet_sync").warning("smart_extract failed for %s: %s", source_name, e)
        return []
    return []
