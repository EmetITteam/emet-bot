import os
import io
import time
import shutil
import pandas as pd
from pptx import Presentation
from docx import Document as DocxDocument
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from pypdf import PdfReader
from dotenv import load_dotenv

# Импорты для ИИ и Базы
from langchain_core.documents import Document
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_community.vectorstores import Chroma
from langchain_text_splitters import RecursiveCharacterTextSplitter

# 1. КОНСТАНТЫ (МАССИВ ПАПОК)
load_dotenv()

FOLDER_KB    = '1RBXHGXOIc2kkSAw-LqzLaRqEE3Ix7L-m'  # Регламенты → db_index_kb_google
FOLDER_COACH = '1KPPBurEoCV_wWzY5HxEtv_TrMI4qXfPa'  # Продукты   → db_index_coach_google
FOLDER_IDS = [FOLDER_KB, FOLDER_COACH]  # для обратной совместимости

SERVICE_ACCOUNT_FILE = 'credentials.json'
SCOPES = ['https://www.googleapis.com/auth/drive.readonly']
GEMINI_KEY = os.getenv("GEMINI_API_KEY")

# 2. ФУНКЦИЯ ПОИСКА ФАЙЛОВ (поддержка корпоративных дисков + отлов ошибок)
def get_all_files_recursive(service, folder_id):
    all_files = []
    page_token = None
    while True:
        query = f"'{folder_id}' in parents and trashed = false"
        try:
            results = service.files().list(
                q=query,
                fields="nextPageToken, files(id, name, mimeType, webViewLink)",
                pageToken=page_token,
                includeItemsFromAllDrives=True,
                supportsAllDrives=True
            ).execute()

            items = results.get('files', [])
            for item in items:
                if item['mimeType'] == 'application/vnd.google-apps.folder':
                    all_files.extend(get_all_files_recursive(service, item['id']))
                else:
                    all_files.append(item)

            page_token = results.get('nextPageToken')
            if not page_token:
                break

        except Exception as e:
            print(f"⚠️ Ошибка доступа к папке {folder_id}. Проверьте, выданы ли права сервисному аккаунту! Ошибка: {e}")
            break

    return all_files

# 3. РАСШИРЕННАЯ ФУНКЦИЯ ЗАГРУЗКИ ТЕКСТА (PDF, Docs, Sheets, Slides, xlsx, pptx)
def download_and_parse_files(service, drive_files):
    all_documents = []
    for file in drive_files:
        try:
            print(f"Обработка: {file['name']} ({file['mimeType']})...")
            text = ""

            # PDF
            if file['mimeType'] == 'application/pdf':
                request = service.files().get_media(fileId=file['id'])
                file_content = io.BytesIO()
                downloader = MediaIoBaseDownload(file_content, request)
                done = False
                while not done: _, done = downloader.next_chunk()
                file_content.seek(0)
                reader = PdfReader(file_content)
                text = "".join([page.extract_text() for page in reader.pages if page.extract_text()])

            # Google Docs
            elif file['mimeType'] == 'application/vnd.google-apps.document':
                request = service.files().export_media(fileId=file['id'], mimeType='text/plain')
                text = request.execute().decode('utf-8')

            # Google Sheets -> CSV
            elif file['mimeType'] == 'application/vnd.google-apps.spreadsheet':
                request = service.files().export_media(fileId=file['id'], mimeType='text/csv')
                text = request.execute().decode('utf-8')

            # Google Slides
            elif file['mimeType'] == 'application/vnd.google-apps.presentation':
                request = service.files().export_media(fileId=file['id'], mimeType='text/plain')
                text = request.execute().decode('utf-8')

            # Excel (.xlsx)
            elif file['mimeType'] == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
                request = service.files().get_media(fileId=file['id'])
                file_content = io.BytesIO()
                downloader = MediaIoBaseDownload(file_content, request)
                done = False
                while not done: _, done = downloader.next_chunk()
                file_content.seek(0)
                df = pd.read_excel(file_content)
                text = df.to_csv(index=False)

            # PowerPoint (.pptx)
            elif file['mimeType'] == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                request = service.files().get_media(fileId=file['id'])
                file_content = io.BytesIO()
                downloader = MediaIoBaseDownload(file_content, request)
                done = False
                while not done: _, done = downloader.next_chunk()
                file_content.seek(0)
                prs = Presentation(file_content)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"

            # Word (.docx)
            elif file['mimeType'] == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                request = service.files().get_media(fileId=file['id'])
                file_content = io.BytesIO()
                downloader = MediaIoBaseDownload(file_content, request)
                done = False
                while not done: _, done = downloader.next_chunk()
                file_content.seek(0)
                doc = DocxDocument(file_content)
                text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])

            # Plain text / CSV (нативные файлы на Drive)
            elif file['mimeType'] in ('text/plain', 'text/csv'):
                request = service.files().get_media(fileId=file['id'])
                text = request.execute().decode('utf-8', errors='replace')

            # Создаем документ для базы
            if text:
                name_lower = file['name'].lower()
                category = "combo" if any(w in name_lower for w in ["протокол", "комбін", "combo"]) else "general"
                metadata = {
                    "source": file['name'],
                    "url": file.get('webViewLink', ''),
                    "category": category
                }
                all_documents.append(Document(page_content=text, metadata=metadata))
        except Exception as e:
            print(f"❌ Ошибка в {file['name']}: {e}")

    return all_documents

# 4. ФУНКЦИЯ СОЗДАНИЯ ВЕКТОРНОЙ БАЗЫ GOOGLE (с защитой от 429)
def build_vector_db(documents):
    embeddings = GoogleGenerativeAIEmbeddings(
        model="models/gemini-embedding-001",
        google_api_key=GEMINI_KEY
    )

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1500, chunk_overlap=200)
    texts = text_splitter.split_documents(documents)

    print(f"Всего фрагментов: {len(texts)}. Начинаю загрузку...")

    vector_db = None
    batch_size = 2  # По 2 куска за раз — самая стабильная скорость для платного старта
    db_path = "data/db_index_google"

    for i in range(0, len(texts), batch_size):
        batch = texts[i:i + batch_size]
        success = False
        while not success:
            try:
                if vector_db is None:
                    vector_db = Chroma.from_documents(
                        documents=batch,
                        embedding=embeddings,
                        persist_directory=db_path
                    )
                else:
                    vector_db.add_documents(batch)
                success = True
                print(f"Прогресс Google: {min(i + batch_size, len(texts))} / {len(texts)}")
                time.sleep(1)  # Небольшая пауза, чтобы не злить API
            except Exception as e:
                err = str(e)
                if "429" in err:
                    print("Превышен лимит (TPM). Сплю 30 секунд...")
                    time.sleep(30)
                elif "503" in err or "UNAVAILABLE" in err:
                    print("Сервис временно недоступен (503). Сплю 15 секунд...")
                    time.sleep(15)
                else:
                    print(f"Ошибка: {e}")
                    raise e

    print(f"✅ База данных для Google успешно собрана в папке: {db_path}")
    return vector_db

# 5. ГЛАВНЫЙ ЗАПУСК — РАЗДЕЛЬНЫЕ БАЗЫ
def _build_to_path_google(documents, db_path):
    """Аналог build_vector_db, но с явным db_path."""
    embeddings = GoogleGenerativeAIEmbeddings(
        model="models/gemini-embedding-001",
        google_api_key=GEMINI_KEY
    )
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1500, chunk_overlap=200)
    texts = text_splitter.split_documents(documents)
    print(f"   Фрагментов: {len(texts)}")
    vector_db = None
    batch_size = 2
    for i in range(0, len(texts), batch_size):
        batch = texts[i:i + batch_size]
        success = False
        while not success:
            try:
                if vector_db is None:
                    vector_db = Chroma.from_documents(documents=batch, embedding=embeddings, persist_directory=db_path)
                else:
                    vector_db.add_documents(batch)
                success = True
                print(f"   Прогресс: {min(i + batch_size, len(texts))} / {len(texts)}")
                time.sleep(1)
            except Exception as e:
                err = str(e)
                if "429" in err:
                    print("   Rate limit, жду 30 сек...")
                    time.sleep(30)
                elif "503" in err or "UNAVAILABLE" in err:
                    print("   Сервис недоступен, жду 15 сек...")
                    time.sleep(15)
                else:
                    raise e


def build_db_for_folder(service, folder_id, db_path, label):
    print(f"\n-> Сканирую папку [{label}]: {folder_id}")
    files = get_all_files_recursive(service, folder_id.strip())
    print(f"   Файлов: {len(files)}")
    docs = download_and_parse_files(service, files)
    if docs:
        if os.path.exists(db_path):
            shutil.rmtree(db_path)
        _build_to_path_google(docs, db_path)
        print(f"   ✅ [{label}] сохранена в {db_path}")
    else:
        print(f"   ⚠️ [{label}] нет документов")


def run_indexing():
    creds = service_account.Credentials.from_service_account_file(SERVICE_ACCOUNT_FILE, scopes=SCOPES)
    service = build('drive', 'v3', credentials=creds)

    print("=== Сборка раздельных баз Google ===")
    build_db_for_folder(service, FOLDER_KB,    "data/db_index_kb_google",    "KB-регламенты")
    build_db_for_folder(service, FOLDER_COACH, "data/db_index_coach_google", "Коуч-продукты")
    print("\n--- ГОТОВО! Базы KB и Коуч сохранены раздельно ---")

if __name__ == "__main__":
    run_indexing()
